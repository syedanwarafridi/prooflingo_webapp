from pptx import Presentation
import pdfplumber
import pandas as pd
from spire.doc import *
from spire.doc.common import *
from .Assistants import comparison_assistant
from docx import Document
from pptx import Presentation
from pptx.util import Inches
import re

# ------> PPTX Document <------#
def extract_text_and_tables_from_pptx(pptx_path):
    presentation = Presentation(pptx_path)
    response_data = []
    for slide in presentation.slides:
        slide_content = ""
        tables = []
        for shape in slide.shapes:
            # Extract text from shapes
            if hasattr(shape, "text"):
                slide_content += shape.text + "\n"
            
            # Extract tables
            if shape.has_table:
                table = shape.table
                table_data = []
                for row in table.rows:
                    row_data = [cell.text.strip() if cell.text else "" for cell in row.cells]
                    table_data.append(row_data)
                tables.append(table_data)
        
        response_data.append({
            "content": slide_content,
            "tables": tables
        })
    
    return response_data

#------> PDf Conversion <-------#
def word_to_pdf(file_path):
    document = Document()
    document.LoadFromFile(file_path)
    
    pdf_path = "pdf_files/Translated_ToPDF.pdf"
    document.SaveToFile(pdf_path, FileFormat.PDF)
    document.Close()
    
    return pdf_path

# Define unwanted substrings
unwanted_substrings = [
    "War\n", 
    "ning: The document was cre\n", 
    "ated with Spire.Doc fo\n", 
    "Note: the cost per impression is $1.50 for the...",
    "on"
]

# Function to clean unwanted substrings from a given string
def clean_text(text, substrings):
    if text is None:
        return text
    for substring in substrings:
        text = text.replace(substring, '')
    return text

# ------> Word Document <-------#
def extract_text_and_tables_word(file_path):
    pdf_path = word_to_pdf(file_path)
    with pdfplumber.open(pdf_path) as pdf:
        all_content_and_tables = []
        
        for i, page in enumerate(pdf.pages):
            # Extract tables first to identify their positions
            tables = page.extract_tables()
            table_bboxes = [table.bbox for table in page.find_tables()]

            # Extract the text
            text = page.extract_text()
            cleaned_text = clean_text(text, unwanted_substrings) if text else ""
            table_positions = []

            if text:
                words = page.extract_words()
                cleaned_words = []
                for word in words:
                    is_table_word = False
                    for bbox in table_bboxes:
                        if (word["x0"] >= bbox[0] and word["x1"] <= bbox[2] and
                                word["top"] >= bbox[1] and word["bottom"] <= bbox[3]):
                            is_table_word = True
                            break
                    if not is_table_word:
                        cleaned_words.append(word)
                    else:
                        if (bbox[1], bbox[3]) not in table_positions:
                            table_positions.append((bbox[1], bbox[3]))

                # Insert <table> tags in the appropriate positions in the text
                cleaned_text_parts = []
                current_position = 0

                for pos in table_positions:
                    table_tag_position = min(
                        range(len(cleaned_words)),
                        key=lambda i: abs((cleaned_words[i]["top"] + cleaned_words[i]["bottom"]) / 2 - (pos[0] + pos[1]) / 2)
                    )

                    # Text before the table
                    text_before_table = " ".join(word["text"] for word in cleaned_words[current_position:table_tag_position])
                    cleaned_text_parts.append(text_before_table)
                    cleaned_text_parts.append("<table>")
                    cleaned_text_parts.append("</table>")  # Add closing tag
                    current_position = table_tag_position

                # Remaining text after the last table
                text_after_last_table = " ".join(word["text"] for word in cleaned_words[current_position:])
                cleaned_text_parts.append(text_after_last_table)
                cleaned_text = " ".join(cleaned_text_parts)
                
                # Process the tables and convert to DataFrames
                page_tables = []
                for table in tables:
                    # Remove unwanted substrings from column names
                    table_columns = [clean_text(col, unwanted_substrings) for col in table[0]]
                    table_rows = [[clean_text(cell, unwanted_substrings) for cell in row] for row in table[1:]]
                    df = pd.DataFrame(table_rows, columns=table_columns)
                    page_tables.append(df)
                
                # Append content and tables for the current page
                all_content_and_tables.append({"content": cleaned_text, "tables": page_tables})
    for i in range(len(all_content_and_tables)):
        all_content_and_tables[i]["content"] = all_content_and_tables[i]["content"].replace("Evaluation Warning: The document was created with Spire.Doc for Python.", "")
    
    return all_content_and_tables

# # ------> Word Document <-------#
# def extract_text_and_tables_word(pdf_path):
#     pdf_path = word_to_pdf(pdf_path)
#     with pdfplumber.open(pdf_path) as pdf:
#         all_content_and_tables = []
        
#         for i, page in enumerate(pdf.pages):
#             # Extract tables first to identify their positions
#             tables = page.extract_tables()
#             table_bboxes = [table.bbox for table in page.find_tables()]

#             # Extract the text
#             text = page.extract_text()
#             cleaned_text = text
#             table_positions = []

#             if text:
#                 words = page.extract_words()
#                 cleaned_words = []
#                 for word in words:
#                     is_table_word = False
#                     for bbox in table_bboxes:
#                         if (word["x0"] >= bbox[0] and word["x1"] <= bbox[2] and
#                                 word["top"] >= bbox[1] and word["bottom"] <= bbox[3]):
#                             is_table_word = True
#                             break
#                     if not is_table_word:
#                         cleaned_words.append(word)
#                     else:
#                         if (bbox[1], bbox[3]) not in table_positions:
#                             table_positions.append((bbox[1], bbox[3]))

#                 # Insert <table> tags in the appropriate positions in the text
#                 cleaned_text_parts = []
#                 current_position = 0

#                 for pos in table_positions:
#                     table_tag_position = min(
#                         range(len(cleaned_words)),
#                         key=lambda i: abs((cleaned_words[i]["top"] + cleaned_words[i]["bottom"]) / 2 - (pos[0] + pos[1]) / 2)
#                     )

#                     # Text before the table
#                     text_before_table = " ".join(word["text"] for word in cleaned_words[current_position:table_tag_position])
#                     cleaned_text_parts.append(text_before_table)
#                     cleaned_text_parts.append("<table>")
#                     cleaned_text_parts.append("</table>")  # Add closing tag
#                     current_position = table_tag_position

#                 # Remaining text after the last table
#                 text_after_last_table = " ".join(word["text"] for word in cleaned_words[current_position:])
#                 cleaned_text_parts.append(text_after_last_table)
#                 cleaned_text = " ".join(cleaned_text_parts)
                
#                 # Process the tables and convert to DataFrames
#                 page_tables = []
#                 for table in tables:
#                     # Remove newline characters from column names
#                     table_columns = [col.replace("\n", " ") for col in table[0]]
#                     table_rows = table[1:]
#                     df = pd.DataFrame(table_rows, columns=table_columns)
#                     page_tables.append(df)
                
#                 # Append content and tables for the current page
#                 all_content_and_tables.append({"content": cleaned_text, "tables": page_tables})
                
#     return all_content_and_tables



# #----> Embedding tables in Text <------#
def embed_tables(text, tables):
    for i, table in enumerate(tables):
        table_str = ""
        for col, values in table.items():
            table_str += f"{col}: {values}\n"
        text = text.replace(f"<table> </table>", table_str, 1)
    return text

# #----> Embedding Tables in segments <------#
def embed_tables_segments(content, tables):
    for i, sentence in enumerate(content):
        for j, table in enumerate(tables):
            table_str = ""
            for col, values in table.items():
                table_str += f"{col}: {values}\n"
            sentence = sentence.replace("<table> </table>", table_str, 1)
        content[i] = sentence
    return content


# #----> Main Comparison Api Running <------#
def process_content(source, target, context, llm_response_history, lang, segment):
    if segment == True:
        source_page = embed_tables_segments(source["content"], source["tables"])
        target_page = embed_tables_segments(target["content"], target["tables"])
    else:
        source_page = embed_tables(source["content"], source["tables"])
        target_page = embed_tables(target["content"], target["tables"])
        
    response = comparison_assistant(source_page, target_page, llm_response_history, context, lang)

    return source_page, target_page, response

# -----> PPTX with Segments <------ #
def extract_text_and_tables_from_pptx_segments(pptx_path):
    presentation = Presentation(pptx_path)
    response_data = []
    for slide in presentation.slides:
        slide_content = []
        tables = []
        for shape in slide.shapes:
            # Extract text from shapes
            if hasattr(shape, "text"):
                text = shape.text.replace("\r", "").replace("\n", " ")
                sentences = text.split(". ")
                slide_content.extend(sentences)
            
            # Extract tables
            if shape.has_table:
                table = shape.table
                table_data = []
                for row in table.rows:
                    row_data = [cell.text.strip() if cell.text else "" for cell in row.cells]
                    table_data.append(row_data)
                tables.append(table_data)
        
        response_data.append({
            "content": slide_content,
            "tables": tables
        })
    
    return response_data

# ------> Word Document with Segments <------ #
# ------> PDF Conversion <-------#
def word_to_pdf(file_path):
    document = Document()
    document.LoadFromFile(file_path)
    
    pdf_path = "pdf_files/Translated_ToPDF.pdf"
    document.SaveToFile(pdf_path, FileFormat.PDF)
    document.Close()
    
    return pdf_path


# ------> Word Document <-------#
def extract_text_and_tables_word_segments(pdf_path):
    pdf_path = word_to_pdf(pdf_path)
    with pdfplumber.open(pdf_path) as pdf:
        all_content_and_tables = []
        
        for i, page in enumerate(pdf.pages):
            # Extract tables first to identify their positions
            tables = page.extract_tables()
            table_bboxes = [table.bbox for table in page.find_tables()]

            # Extract the text
            text = page.extract_text()
            cleaned_text = text
            table_positions = []

            if text:
                words = page.extract_words()
                cleaned_words = []
                for word in words:
                    is_table_word = False
                    for bbox in table_bboxes:
                        if (word["x0"] >= bbox[0] and word["x1"] <= bbox[2] and
                                word["top"] >= bbox[1] and word["bottom"] <= bbox[3]):
                            is_table_word = True
                            break
                    if not is_table_word:
                        cleaned_words.append(word)
                    else:
                        if (bbox[1], bbox[3]) not in table_positions:
                            table_positions.append((bbox[1], bbox[3]))

                # Insert <table> tags in the appropriate positions in the text
                cleaned_text_parts = []
                current_position = 0

                for pos in table_positions:
                    table_tag_position = min(
                        range(len(cleaned_words)),
                        key=lambda i: abs((cleaned_words[i]["top"] + cleaned_words[i]["bottom"]) / 2 - (pos[0] + pos[1]) / 2)
                    )

                    # Text before the table
                    text_before_table = " ".join(word["text"] for word in cleaned_words[current_position:table_tag_position])
                    cleaned_text_parts.append(text_before_table)
                    cleaned_text_parts.append("<table>")
                    cleaned_text_parts.append("</table>")  # Add closing tag
                    current_position = table_tag_position

                # Remaining text after the last table
                text_after_last_table = " ".join(word["text"] for word in cleaned_words[current_position:])
                cleaned_text_parts.append(text_after_last_table)
                cleaned_text = " ".join(cleaned_text_parts)
                
                # Split cleaned text into sentences
                sentences = re.split(r'(?<=[.!?]) +', cleaned_text)
                
                # Process the tables and convert to DataFrames
                page_tables = []
                for table in tables:
                    if table:
                        # Remove newline characters from column names
                        table_columns = [col.replace("\n", " ") if col else "" for col in table[0]]
                        table_rows = [[cell if cell is not None else "" for cell in row] for row in table[1:]]
                        df = pd.DataFrame(table_rows, columns=table_columns)
                        page_tables.append(df)
                
                # Append content and tables for the current page
                all_content_and_tables.append({"content": sentences, "tables": page_tables})
                
    return all_content_and_tables

#------------> Auto Save <----------#
def process_docx(updated_file_path, comparison_data):
    if os.path.exists(updated_file_path):
        document = Document(updated_file_path)
    else:
        document = Document()

    document.add_page_break()
    document.add_paragraph(comparison_data.get('content', ''))

    for table_data in comparison_data.get('tables', []):
        table = document.add_table(rows=1, cols=len(table_data))
        hdr_cells = table.rows[0].cells
        for i, key in enumerate(table_data):
            hdr_cells[i].text = key

        for i in range(len(next(iter(table_data.values())))):
            row_cells = table.add_row().cells
            for j, key in enumerate(table_data):
                row_cells[j].text = table_data[key][i]

    document.save(updated_file_path)

def process_pptx(updated_file_path, comparison_data):
    if os.path.exists(updated_file_path):
        presentation = Presentation(updated_file_path)
    else:
        presentation = Presentation()

    slide_layout = presentation.slide_layouts[5]  # Assuming layout 5 for blank slide
    slide = presentation.slides.add_slide(slide_layout)
    text_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(8.5), Inches(1.5))
    text_frame = text_box.text_frame
    text_frame.text = comparison_data.get('content', '')

    for table_data in comparison_data.get('tables', []):
        rows = len(next(iter(table_data.values()))) + 1
        cols = len(table_data)
        table_shape = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(2), Inches(9), Inches(5)).table
        
        for j, key in enumerate(table_data):
            table_shape.cell(0, j).text = key

        for i in range(1, rows):
            for j, key in enumerate(table_data):
                table_shape.cell(i, j).text = table_data[key][i-1]

    presentation.save(updated_file_path)
    
#---------> Bilingual File Reader <-------------#
import xml.etree.ElementTree as ET
import re

def parse_bilingual(file_path):
    tree = ET.parse(file_path)
    root = tree.getroot()
    
    namespaces = {'xliff': 'urn:oasis:names:tc:xliff:document:1.2'}
    source_texts = []
    target_texts = []

    def strip_special_tags(text):
        # Strip <fontFormat> tags and their content
        text = re.sub(r'<fontFormat.?>(.?)<\/fontFormat>', r'\1', text, flags=re.DOTALL)
        # Strip <fieldStart> and <fieldEnd> tags
        text = re.sub(r'<fieldStart.*?>', '', text)
        text = re.sub(r'<fieldEnd.*?>', '', text)
        return text.strip()  # Strip leading and trailing whitespace

    for trans_unit in root.findall(".//xliff:trans-unit", namespaces):
        source = trans_unit.find("xliff:source", namespaces)
        target = trans_unit.find("xliff:target", namespaces)
        
        source_text = ''.join(source.itertext()).strip() if source is not None else ''
        target_text = ''.join(target.itertext()).strip() if target is not None else ''

        source_texts.append(strip_special_tags(source_text))
        target_texts.append(strip_special_tags(target_text))

    # Join all source texts and target texts with newline characters
    all_source_text = '\n'.join(source_texts)
    all_target_text = '\n'.join(target_texts)

    return all_source_text, all_target_text

# # Example file path, replace with your actual file path
# file_path = 'cleint.txlf'  # Replace with your actual file path

# source_text, target_text = parse_bilingual(file_path)

# print("Source Texts:")
# print(source_text)

# print("\nTarget Texts:")
# print(target_text)
########################## -------> Content Extraction Using New Flow of XML <-------- ############################
import os
from docx_utils.flatten import opc_to_flat_opc

def docx_to_xml(file_path):
    dir_path = os.path.join("xml_files")
    if not os.path.exists(dir_path):
        os.makedirs(dir_path)

    # Extract the filename without the extension
    name = os.path.splitext(os.path.basename(file_path))[0]

    out_file_path = os.path.join(dir_path, f"{name}.xml")
    opc_to_flat_opc(file_path, out_file_path)

    return out_file_path

# docx_to_xml("new-source.docx")

import xml.etree.ElementTree as ET
import pandas as pd

def extract_text_and_tables_from_xml(file_path):
    xml_file_path = docx_to_xml(file_path)
    # Load and parse the XML file
    tree = ET.parse(xml_file_path)
    root = tree.getroot()

    # Namespace dictionary to handle XML namespaces in the document
    namespaces = {
        'w': 'http://schemas.openxmlformats.org/wordprocessingml/2006/main'
    }

    # Find all body elements
    body_elements = root.findall('.//w:body', namespaces)

    pages = []
    current_page_content = []

    def extract_run_text(run):
        # Extract text from a run
        text_elements = run.findall('.//w:t', namespaces)
        run_text = ''.join([text_elem.text for text_elem in text_elements if text_elem.text])
        return run_text

    def extract_paragraph_text(paragraph):
        # Extract text from a paragraph, separating lines with \n
        runs = paragraph.findall('.//w:r', namespaces)
        lines = [extract_run_text(run) for run in runs]
        return '\n'.join(lines)

    def extract_table_data(table):
        # Extract table data into a list of lists, separating cell lines with \n
        table_data = []
        for row in table.findall('.//w:tr', namespaces):
            row_data = []
            for cell in row.findall('.//w:tc', namespaces):
                cell_text = '\n'.join(cell.itertext()).strip()
                cell_metadata = extract_metadata(cell)
                row_data.append({
                    'text': cell_text,
                    'metadata': cell_metadata
                })
            table_data.append(row_data)
        return table_data

    def extract_metadata(elem):
        # Extract metadata such as style, alignment, font, and font size
        metadata = {}
        pPr = elem.find('.//w:pPr', namespaces)
        if pPr is not None:
            style = pPr.find('.//w:pStyle', namespaces)
            if style is not None:
                metadata['style'] = style.attrib.get('{http://schemas.openxmlformats.org/wordprocessingml/2006/main}val', '')
            alignment = pPr.find('.//w:jc', namespaces)
            if alignment is not None:
                metadata['alignment'] = alignment.attrib.get('{http://schemas.openxmlformats.org/wordprocessingml/2006/main}val', '')

        rPr = elem.find('.//w:rPr', namespaces)
        if rPr is not None:
            font = rPr.find('.//w:rFonts', namespaces)
            if font is not None:
                metadata['font'] = font.attrib.get('{http://schemas.openxmlformats.org/wordprocessingml/2006/main}ascii', '')
            font_size = rPr.find('.//w:sz', namespaces)
            if font_size is not None:
                fontsize = int(font_size.attrib.get('{http://schemas.openxmlformats.org/wordprocessingml/2006/main}val', '')) // 2
                metadata['font_size'] = str(fontsize)
            if rPr.find('.//w:b', namespaces) is not None:
                metadata['bold'] = True
            if rPr.find('.//w:i', namespaces) is not None:
                metadata['italic'] = True
            if rPr.find('.//w:u', namespaces) is not None:
                metadata['underline'] = True
        return metadata

    def extract_paragraph_metadata(paragraph):
        # Extract metadata for each run within a paragraph
        runs = paragraph.findall('.//w:r', namespaces)
        metadata_list = []
        for run in runs:
            run_text = extract_run_text(run)
            run_metadata = extract_metadata(run)
            metadata_list.append({
                'text': run_text,
                'metadata': run_metadata
            })
        return metadata_list

    for body in body_elements:
        for elem in body:
            if elem.tag == '{http://schemas.openxmlformats.org/wordprocessingml/2006/main}p':
                if elem.find('.//w:r//w:lastRenderedPageBreak', namespaces) is not None:
                    # Add the current content to pages and start a new page
                    pages.append(current_page_content)
                    current_page_content = []

                para_text = extract_paragraph_text(elem)
                para_metadata = extract_paragraph_metadata(elem)
                paragraph_style = extract_metadata(elem).get('style', '')
                current_page_content.append({
                    'type': 'paragraph',
                    'text': para_text,
                    'metadata': para_metadata,
                    'style': paragraph_style
                })

            elif elem.tag == '{http://schemas.openxmlformats.org/wordprocessingml/2006/main}tbl':  # Table element
                table_data = extract_table_data(elem)
                table_metadata = extract_metadata(elem)
                current_page_content.append({
                    'type': 'table',
                    'data': table_data,
                    'metadata': table_metadata
                })

                # Add the current content to pages and reset for next content
                pages.append(current_page_content)
                current_page_content = []

    # Add any remaining content as the last page
    if current_page_content:
        pages.append(current_page_content)
    return pages

#------> Converting into Lists <-------#
def extract_paragraphs_and_tables_into_lists(pages):
    all_paragraphs_and_tables_per_page = []

    for page in pages:
        page_content = []

        for element in page:
            if element['type'] == 'paragraph':
                page_content.append(element['text'].strip())

            elif element['type'] == 'table':
                table_content = []

                for row in element['data']:
                    row_text = [cell['text'].strip() for cell in row if cell['text'].strip() != '']
                    table_content.append(row_text)

                page_content.append(table_content)  # Append list of table rows

        all_paragraphs_and_tables_per_page.append(page_content)

    return all_paragraphs_and_tables_per_page